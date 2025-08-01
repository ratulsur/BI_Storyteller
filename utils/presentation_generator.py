import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from PIL import Image
import io
import base64
import os
from groq import Groq
import json

class PresentationGenerator:
    def __init__(self):
        self.groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))
        self.templates = {
            "Modern Business": {
                "background_color": (45, 55, 72),
                "title_color": (255, 255, 255),
                "text_color": (226, 232, 240),
                "accent_color": (56, 178, 172)
            },
            "Clean Minimal": {
                "background_color": (255, 255, 255),
                "title_color": (45, 55, 72),
                "text_color": (74, 85, 104),
                "accent_color": (49, 130, 206)
            },
            "Corporate Blue": {
                "background_color": (23, 37, 84),
                "title_color": (255, 255, 255),
                "text_color": (203, 213, 224),
                "accent_color": (66, 153, 225)
            },
            "Warm Orange": {
                "background_color": (255, 248, 240),
                "title_color": (194, 65, 12),
                "text_color": (124, 45, 18),
                "accent_color": (251, 146, 60)
            },
            "Professional Green": {
                "background_color": (240, 253, 244),
                "title_color": (22, 101, 52),
                "text_color": (20, 83, 45),
                "accent_color": (72, 187, 120)
            }
        }
    
    def create_comprehensive_presentation(self, session_state, template_name="Modern Business"):
        """Create a comprehensive presentation from all available session data"""
        prs = Presentation()
        template = self.templates[template_name]
        
        # Title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = title_slide.shapes.title
        subtitle = title_slide.placeholders[1]
        
        business_problem = session_state.get('business_problem', '')
        title_text = "Business Intelligence Report"
        if business_problem:
            # Extract key topic from business problem
            topic_keywords = business_problem.split()[:5]
            title_text = f"Analysis: {' '.join(topic_keywords)}..."
        
        if title and hasattr(title, 'text_frame'):
            title.text = title_text
        if subtitle and hasattr(subtitle, 'text_frame'):
            subtitle.text = "AI-Powered Data Analysis Results"
        
        self._apply_template_styling(title_slide, template)
        
        # Business Problem slide
        if business_problem:
            problem_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = problem_slide.shapes.title
            content = problem_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Business Problem"
            if content and hasattr(content, 'text_frame'):
                content.text = business_problem
            
            self._apply_template_styling(problem_slide, template)
        
        # Data Overview slide
        processed_data = session_state.get('processed_data')
        if processed_data is not None:
            data_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = data_slide.shapes.title
            content = data_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Dataset Overview"
            if content and hasattr(content, 'text_frame'):
                data_info = f"""Dataset Statistics:
• Total Records: {processed_data.shape[0]:,}
• Variables Analyzed: {processed_data.shape[1]}
• Data Quality: {((1 - processed_data.isnull().sum().sum() / (processed_data.shape[0] * processed_data.shape[1])) * 100):.1f}%
• Key Variables: {', '.join(processed_data.columns[:5].tolist())}"""
                content.text = data_info
            
            self._apply_template_styling(data_slide, template)
            
            # Data Quality Details slide
            quality_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = quality_slide.shapes.title
            content = quality_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Data Quality Assessment"
            if content and hasattr(content, 'text_frame'):
                missing_data = processed_data.isnull().sum()
                top_missing = missing_data.head(5)
                
                quality_info = f"""Data Completeness Analysis:

• Overall Data Completeness: {((1 - processed_data.isnull().sum().sum() / (processed_data.shape[0] * processed_data.shape[1])) * 100):.1f}%
• Records with Complete Data: {len(processed_data.dropna()):,}
• Variables with Missing Data: {(missing_data > 0).sum()}

Top Variables by Completeness:
{chr(10).join([f"• {col}: {((len(processed_data) - missing) / len(processed_data) * 100):.1f}% complete" for col, missing in top_missing.items()])}"""
                content.text = quality_info
            
            self._apply_template_styling(quality_slide, template)
            
            # Variable Analysis slide
            if processed_data.shape[1] > 5:
                var_slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = var_slide.shapes.title
                content = var_slide.placeholders[1]
                
                if title and hasattr(title, 'text_frame'):
                    title.text = "Variable Analysis"
                if content and hasattr(content, 'text_frame'):
                    numeric_cols = processed_data.select_dtypes(include=['number']).columns
                    categorical_cols = processed_data.select_dtypes(include=['object']).columns
                    
                    var_info = f"""Variable Type Distribution:

• Numeric Variables: {len(numeric_cols)} 
  {', '.join(numeric_cols[:8].tolist())}{'...' if len(numeric_cols) > 8 else ''}

• Categorical Variables: {len(categorical_cols)}
  {', '.join(categorical_cols[:8].tolist())}{'...' if len(categorical_cols) > 8 else ''}

• Total Unique Values Across Dataset: {processed_data.nunique().sum():,}
• Average Unique Values per Variable: {processed_data.nunique().mean():.1f}"""
                    content.text = var_info
                
                self._apply_template_styling(var_slide, template)
        
        # Analysis Results slides
        analysis_data = session_state.get('analysis_results', {})
        
        # Executive Summary
        if analysis_data.get('summary'):
            summary_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = summary_slide.shapes.title
            content = summary_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Executive Summary"
            if content and hasattr(content, 'text_frame'):
                content.text = analysis_data['summary']
            
            self._apply_template_styling(summary_slide, template)
        
        # Key Insights - Split into multiple slides if needed
        if analysis_data.get('insights'):
            insights_list = analysis_data['insights']
            if isinstance(insights_list, list):
                # Split insights into chunks of 4-5 per slide
                insights_chunks = [insights_list[i:i+4] for i in range(0, len(insights_list), 4)]
                
                for i, chunk in enumerate(insights_chunks):
                    insights_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = insights_slide.shapes.title
                    content = insights_slide.placeholders[1]
                    
                    if title and hasattr(title, 'text_frame'):
                        if len(insights_chunks) > 1:
                            title.text = f"Key Insights - Part {i+1}"
                        else:
                            title.text = "Key Insights"
                    if content and hasattr(content, 'text_frame'):
                        insights_text = "\n\n".join([f"• {insight}" for insight in chunk])
                        content.text = insights_text
                    
                    self._apply_template_styling(insights_slide, template)
            else:
                insights_slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = insights_slide.shapes.title
                content = insights_slide.placeholders[1]
                
                if title and hasattr(title, 'text_frame'):
                    title.text = "Key Insights"
                if content and hasattr(content, 'text_frame'):
                    content.text = str(insights_list)
                
                self._apply_template_styling(insights_slide, template)
        
        # Sentiment Analysis
        if analysis_data.get('sentiment_analysis'):
            sentiment_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = sentiment_slide.shapes.title
            content = sentiment_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Sentiment Analysis"
            if content and hasattr(content, 'text_frame'):
                sentiment_data = analysis_data['sentiment_analysis']
                if isinstance(sentiment_data, dict):
                    sentiment_text = ""
                    for key, value in sentiment_data.items():
                        sentiment_text += f"• {key}: {value}\n"
                else:
                    sentiment_text = str(sentiment_data)
                content.text = sentiment_text
            
            self._apply_template_styling(sentiment_slide, template)
        
        # Predictions
        if analysis_data.get('predictions') or analysis_data.get('forecast'):
            pred_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = pred_slide.shapes.title
            content = pred_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Predictions & Forecasts"
            if content and hasattr(content, 'text_frame'):
                pred_data = analysis_data.get('predictions') or analysis_data.get('forecast')
                if isinstance(pred_data, dict):
                    pred_text = ""
                    for key, value in pred_data.items():
                        pred_text += f"• {key}: {value}\n"
                elif isinstance(pred_data, list):
                    pred_text = "\n".join([f"• {pred}" for pred in pred_data])
                else:
                    pred_text = str(pred_data)
                content.text = pred_text
            
            self._apply_template_styling(pred_slide, template)
        
        # Dashboard Insights
        dashboard_insights = session_state.get('dashboard_insights')
        if dashboard_insights:
            dashboard_slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = dashboard_slide.shapes.title
            content = dashboard_slide.placeholders[1]
            
            if title and hasattr(title, 'text_frame'):
                title.text = "Dashboard Insights"
            if content and hasattr(content, 'text_frame'):
                if isinstance(dashboard_insights, dict):
                    dashboard_text = ""
                    for key, value in dashboard_insights.items():
                        dashboard_text += f"• {key}: {value}\n"
                else:
                    dashboard_text = str(dashboard_insights)
                content.text = dashboard_text
            
            self._apply_template_styling(dashboard_slide, template)
        
        # Latest Chat Conversation - Multiple slides for detailed conversation
        chat_history = session_state.get('chat_history', [])
        if chat_history:
            # Get the last 4 messages for more context
            last_messages = chat_history[-4:] if len(chat_history) >= 4 else chat_history
            
            if last_messages:
                # Create separate slides for user questions and AI responses
                user_questions = [msg for msg in last_messages if msg.get('role') == 'user']
                ai_responses = [msg for msg in last_messages if msg.get('role') == 'assistant']
                
                if user_questions:
                    q_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = q_slide.shapes.title
                    content = q_slide.placeholders[1]
                    
                    if title and hasattr(title, 'text_frame'):
                        title.text = "Recent User Questions"
                    if content and hasattr(content, 'text_frame'):
                        questions_text = ""
                        for i, msg in enumerate(user_questions[-2:], 1):
                            questions_text += f"{i}. {msg.get('content', '')[:250]}...\n\n"
                        content.text = questions_text
                    
                    self._apply_template_styling(q_slide, template)
                
                if ai_responses:
                    ai_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = ai_slide.shapes.title
                    content = ai_slide.placeholders[1]
                    
                    if title and hasattr(title, 'text_frame'):
                        title.text = "AI Analysis & Insights"
                    if content and hasattr(content, 'text_frame'):
                        ai_text = "Latest AI Analysis:\n\n"
                        latest_response = ai_responses[-1] if ai_responses else None
                        if latest_response:
                            ai_text += latest_response.get('content', '')[:500] + "..."
                        content.text = ai_text
                    
                    self._apply_template_styling(ai_slide, template)
        
        # Recommendations - Split into multiple slides
        if analysis_data.get('recommendations'):
            recommendations = analysis_data['recommendations']
            if isinstance(recommendations, list):
                # Split recommendations into chunks of 3-4 per slide
                rec_chunks = [recommendations[i:i+3] for i in range(0, len(recommendations), 3)]
                
                for i, chunk in enumerate(rec_chunks):
                    rec_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = rec_slide.shapes.title
                    content = rec_slide.placeholders[1]
                    
                    if title and hasattr(title, 'text_frame'):
                        if len(rec_chunks) > 1:
                            title.text = f"Recommendations - Part {i+1}"
                        else:
                            title.text = "Recommendations"
                    if content and hasattr(content, 'text_frame'):
                        rec_text = "\n\n".join([f"• {rec}" for rec in chunk])
                        content.text = rec_text
                    
                    self._apply_template_styling(rec_slide, template)
            else:
                rec_slide = prs.slides.add_slide(prs.slide_layouts[1])
                title = rec_slide.shapes.title
                content = rec_slide.placeholders[1]
                
                if title and hasattr(title, 'text_frame'):
                    title.text = "Recommendations"
                if content and hasattr(content, 'text_frame'):
                    content.text = str(recommendations)
                
                self._apply_template_styling(rec_slide, template)
        
        # Methodology slide
        method_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = method_slide.shapes.title
        content = method_slide.placeholders[1]
        
        if title and hasattr(title, 'text_frame'):
            title.text = "Analysis Methodology"
        if content and hasattr(content, 'text_frame'):
            method_text = """AI-Powered Analysis Approach:

• Business Problem Definition & Variable Extraction
• Custom Questionnaire Generation using LLM
• Automated Data Collection & Quality Assessment
• Advanced Statistical Analysis & Pattern Recognition
• Sentiment Analysis & Predictive Modeling
• Interactive Dashboard Creation
• AI-Driven Insight Generation"""
            content.text = method_text
        
        self._apply_template_styling(method_slide, template)
        
        # Implementation Timeline slide
        timeline_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = timeline_slide.shapes.title
        content = timeline_slide.placeholders[1]
        
        if title and hasattr(title, 'text_frame'):
            title.text = "Implementation Timeline"
        if content and hasattr(content, 'text_frame'):
            timeline_text = """Recommended Implementation Schedule:

Week 1-2: Immediate Actions
• Implement high-priority recommendations
• Set up monitoring systems

Week 3-4: Process Improvements
• Deploy identified solutions
• Train relevant personnel

Month 2-3: Optimization
• Monitor results and adjust strategies
• Collect additional data for refinement

Ongoing: Continuous Improvement
• Regular analysis updates
• Performance tracking"""
            content.text = timeline_text
        
        self._apply_template_styling(timeline_slide, template)
        
        # Next Steps slide
        next_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = next_slide.shapes.title
        content = next_slide.placeholders[1]
        
        if title and hasattr(title, 'text_frame'):
            title.text = "Next Steps & Actions"
        if content and hasattr(content, 'text_frame'):
            next_text = """Immediate Actions Required:

• Review and approve recommended strategies
• Assign ownership for each action item
• Establish success metrics and KPIs
• Schedule follow-up analysis sessions
• Implement data collection improvements
• Prepare stakeholder communication plan"""
            content.text = next_text
        
        self._apply_template_styling(next_slide, template)
        
        # Conclusion slide
        conclusion_slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = conclusion_slide.shapes.title
        content = conclusion_slide.placeholders[1]
        
        if title and hasattr(title, 'text_frame'):
            title.text = "Conclusion"
        if content and hasattr(content, 'text_frame'):
            conclusion_text = """Key Takeaways:

• Data-driven insights provide clear direction
• AI analysis reveals previously hidden patterns
• Recommended actions are prioritized by impact
• Continuous monitoring ensures sustained success

Thank you for using BI StoryTeller AI Platform
For questions or support: Continue using the Chat feature"""
            content.text = conclusion_text
        
        self._apply_template_styling(conclusion_slide, template)
        
        return prs
    
    def create_presentation_from_analysis(self, analysis_data, template_name="Modern Business"):
        """Legacy method for backward compatibility"""
        return self.create_comprehensive_presentation({'analysis_results': analysis_data}, template_name)
    
    def _apply_template_styling(self, slide, template):
        """Apply template styling to slide"""
        # Set background color
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*template["background_color"])
        
        # Style title and content
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if shape == slide.shapes.title:
                            run.font.color.rgb = RGBColor(*template["title_color"])
                            run.font.size = Pt(44)
                            run.font.bold = True
                        else:
                            run.font.color.rgb = RGBColor(*template["text_color"])
                            run.font.size = Pt(24)
    
    def add_image_to_slide(self, prs, slide_index, image_path, left=Inches(1), top=Inches(2), width=Inches(8)):
        """Add an image to a specific slide"""
        if slide_index < len(prs.slides):
            slide = prs.slides[slide_index]
            slide.shapes.add_picture(image_path, left, top, width=width)
            return True
        return False
    
    def generate_image_description(self, slide_content, context="business presentation"):
        """Generate image description using Groq"""
        try:
            prompt = f"""Based on this slide content: "{slide_content}"
            
            Generate a professional image description for a {context} that would complement this content.
            The description should be suitable for AI image generation and focus on:
            - Professional business imagery
            - Clean, modern aesthetic
            - Relevant visual metaphors
            - Appropriate colors and composition
            
            Respond with only the image description, no additional text."""
            
            response = self.groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.7,
                max_tokens=150
            )
            
            content = response.choices[0].message.content
            return content.strip() if content else f"Professional business illustration related to: {slide_content[:100]}"
        except Exception as e:
            return f"Professional business illustration related to: {slide_content[:100]}"
    
    def create_custom_slide(self, prs, title, content, template_name="Modern Business"):
        """Create a custom slide with title and content"""
        template = self.templates[template_name]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        slide_title = slide.shapes.title
        slide_content = slide.placeholders[1]
        
        if slide_title and hasattr(slide_title, 'text_frame'):
            slide_title.text = title
        if slide_content and hasattr(slide_content, 'text_frame'):
            slide_content.text = content
        
        self._apply_template_styling(slide, template)
        return slide
    
    def save_presentation(self, prs, filename):
        """Save presentation to file"""
        filepath = f"presentations/{filename}"
        os.makedirs("presentations", exist_ok=True)
        prs.save(filepath)
        return filepath
    
    def get_slide_preview_text(self, prs):
        """Get text preview of all slides"""
        slides_preview = []
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text_frame") and shape.text_frame.text:
                    slide_text.append(shape.text_frame.text)
            
            slides_preview.append({
                "slide_number": i + 1,
                "content": "\n".join(slide_text)
            })
        
        return slides_preview

class ImageGenerator:
    def __init__(self):
        self.groq_client = Groq(api_key=os.environ.get("GROQ_API_KEY"))
    
    def generate_svg_image(self, prompt):
        """Generate SVG image description using Groq for business presentations"""
        try:
            svg_prompt = f"""Create a simple, professional SVG image for: {prompt}
            
            Generate SVG code that creates a clean, business-appropriate illustration.
            Requirements:
            - Use professional colors (blues, grays, greens)
            - Simple geometric shapes and icons
            - Minimal design suitable for presentations
            - Size: 400x300 viewBox
            - Include relevant business icons or charts
            
            Respond with only the SVG code, starting with <svg> and ending with </svg>."""
            
            response = self.groq_client.chat.completions.create(
                model="llama-3.3-70b-versatile",
                messages=[{"role": "user", "content": svg_prompt}],
                temperature=0.3,
                max_tokens=1000
            )
            
            content = response.choices[0].message.content
            svg_content = content.strip() if content else ""
            
            # Extract SVG if wrapped in code blocks
            if "```" in svg_content:
                svg_content = svg_content.split("```")[1]
                if svg_content.startswith("svg"):
                    svg_content = svg_content[3:]
            
            # Ensure it starts with <svg
            if not svg_content.strip().startswith("<svg"):
                return self._create_fallback_svg(prompt)
            
            return svg_content
            
        except Exception as e:
            return self._create_fallback_svg(prompt)
    
    def _create_fallback_svg(self, prompt):
        """Create a simple fallback SVG"""
        return f'''<svg viewBox="0 0 400 300" xmlns="http://www.w3.org/2000/svg">
            <rect width="400" height="300" fill="#f8f9fa" stroke="#dee2e6" stroke-width="2"/>
            <circle cx="200" cy="150" r="60" fill="#007bff" opacity="0.8"/>
            <text x="200" y="220" text-anchor="middle" font-family="Arial" font-size="16" fill="#495057">
                {prompt[:30]}...
            </text>
        </svg>'''
    
    def save_svg_as_image(self, svg_content, filename):
        """Save SVG content as image file"""
        os.makedirs("generated_images", exist_ok=True)
        filepath = f"generated_images/{filename}.svg"
        
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(svg_content)
        
        return filepath