import streamlit as st
from utils.presentation_generator import PresentationGenerator, ImageGenerator
import pandas as pd
import os
from pptx import Presentation
import tempfile

st.set_page_config(
    page_title="Presentation Builder - BI StoryTeller",
    layout="wide"
)

st.title("Presentation Builder")
st.markdown("Create professional PowerPoint presentations with AI-generated content and images.")

# Initialize generators
if 'presentation_generator' not in st.session_state:
    st.session_state.presentation_generator = PresentationGenerator()

if 'image_generator' not in st.session_state:
    st.session_state.image_generator = ImageGenerator()

if 'current_presentation' not in st.session_state:
    st.session_state.current_presentation = None

if 'generated_images' not in st.session_state:
    st.session_state.generated_images = []

# Get available data
analysis_data = st.session_state.get('analysis_results')
processed_data = st.session_state.get('processed_data')
business_problem = st.session_state.get('business_problem')
chat_history = st.session_state.get('chat_history')
dashboard_insights = st.session_state.get('dashboard_insights')

# Show data availability status
data_sources_count = sum([
    bool(business_problem),
    bool(processed_data is not None),
    bool(analysis_data),
    bool(chat_history),
    bool(dashboard_insights)
])

if data_sources_count == 0:
    st.warning("No data available for presentation generation. Please complete the workflow steps first.")
    col1, col2 = st.columns(2)
    with col1:
        if st.button("Go to Analysis Page"):
            st.switch_page("pages/5_Analysis.py")
    with col2:
        if st.button("Start with Variables"):
            st.switch_page("pages/1_Variables.py")
else:
    st.success(f"Data from {data_sources_count} sources available for presentation generation!")

# Presentation Creation Section
st.header("Create Presentation")

tab1, tab2, tab3 = st.tabs(["Quick Generate", "Custom Builder", "AI Image Studio"])

with tab1:
    st.subheader("Generate Comprehensive Presentation")
    
    # Show available data sources
    available_data = []
    if st.session_state.get('business_problem'):
        available_data.append("Business Problem Definition")
    if processed_data is not None:
        available_data.append(f"Dataset ({processed_data.shape[0]:,} records)")
    if analysis_data:
        if analysis_data.get('summary'):
            available_data.append("Executive Summary")
        if analysis_data.get('insights'):
            available_data.append(f"Key Insights ({len(analysis_data['insights'])} items)")
        if analysis_data.get('sentiment_analysis'):
            available_data.append("Sentiment Analysis")
        if analysis_data.get('predictions') or analysis_data.get('forecast'):
            available_data.append("Predictions & Forecasts")
        if analysis_data.get('recommendations'):
            available_data.append(f"Recommendations ({len(analysis_data['recommendations'])} items)")
    if st.session_state.get('dashboard_insights'):
        available_data.append("Dashboard Insights")
    if st.session_state.get('chat_history'):
        available_data.append("Latest AI Chat Conversation")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.markdown("**Available Data for Presentation:**")
        if available_data:
            for data_source in available_data:
                st.markdown(f"✓ {data_source}")
        else:
            st.warning("No analysis data available. Complete the workflow to generate comprehensive presentations.")
    
    with col2:
        template_choice = st.selectbox(
            "Choose Template",
            options=list(st.session_state.presentation_generator.templates.keys()),
            key="quick_template"
        )
        
        if st.button("Generate Complete Presentation", type="primary", disabled=not available_data):
            with st.spinner("Creating comprehensive presentation..."):
                prs = st.session_state.presentation_generator.create_comprehensive_presentation(
                    st.session_state, template_choice
                )
                st.session_state.current_presentation = prs
                st.success(f"Comprehensive presentation created with {len(available_data)} data sections!")
                st.balloons()
                st.rerun()
    
    # Data preview section
    if available_data:
        with st.expander("Preview Available Data"):
            if st.session_state.get('business_problem'):
                st.markdown("**Business Problem:**")
                st.write(st.session_state.business_problem[:200] + "...")
            
            if analysis_data and analysis_data.get('summary'):
                st.markdown("**Executive Summary:**")
                st.write(analysis_data['summary'][:200] + "...")
            
            if st.session_state.get('chat_history'):
                st.markdown("**Latest Chat:**")
                recent_chat = st.session_state.chat_history[-1] if st.session_state.chat_history else None
                if recent_chat:
                    st.write(f"{recent_chat.get('role', 'Unknown')}: {recent_chat.get('content', '')[:150]}...")
    
    st.markdown("---")
    st.info("💡 **Tip:** Complete more workflow steps (Analysis, Dashboard, Chat) to generate richer presentations with comprehensive insights.")

with tab2:
    st.subheader("Custom Slide Builder")
    
    # Template selection
    template_choice = st.selectbox(
        "Select Template",
        options=list(st.session_state.presentation_generator.templates.keys()),
        key="custom_template"
    )
    
    # Show template preview
    template_info = st.session_state.presentation_generator.templates[template_choice]
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown("**Background**")
        bg_color = f"rgb({template_info['background_color'][0]}, {template_info['background_color'][1]}, {template_info['background_color'][2]})"
        st.markdown(f'<div style="background-color: {bg_color}; height: 50px; border-radius: 5px;"></div>', unsafe_allow_html=True)
    
    with col2:
        st.markdown("**Title Color**")
        title_color = f"rgb({template_info['title_color'][0]}, {template_info['title_color'][1]}, {template_info['title_color'][2]})"
        st.markdown(f'<div style="background-color: {title_color}; height: 50px; border-radius: 5px;"></div>', unsafe_allow_html=True)
    
    with col3:
        st.markdown("**Accent Color**")
        accent_color = f"rgb({template_info['accent_color'][0]}, {template_info['accent_color'][1]}, {template_info['accent_color'][2]})"
        st.markdown(f'<div style="background-color: {accent_color}; height: 50px; border-radius: 5px;"></div>', unsafe_allow_html=True)
    
    # Create new presentation
    if st.button("Create New Presentation"):
        prs = Presentation()
        # Add title slide
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = "New Presentation"
        title_slide.placeholders[1].text = "Created with BI StoryTeller"
        
        st.session_state.current_presentation = prs
        st.success("New presentation created!")
        st.rerun()
    
    # Add slides to existing presentation
    if st.session_state.current_presentation:
        st.markdown("---")
        st.subheader("Add New Slide")
        
        slide_title = st.text_input("Slide Title", placeholder="Enter slide title...")
        slide_content = st.text_area("Slide Content", placeholder="Enter slide content...", height=150)
        
        if st.button("Add Slide") and slide_title:
            st.session_state.presentation_generator.create_custom_slide(
                st.session_state.current_presentation,
                slide_title,
                slide_content,
                template_choice
            )
            st.success(f"Added slide: {slide_title}")
            st.rerun()

with tab3:
    st.subheader("AI Image Studio")
    st.markdown("Generate custom images for your presentation using AI.")
    
    # Image generation chat interface
    st.markdown("### Image Generation Chat")
    
    # Chat history
    if 'image_chat_history' not in st.session_state:
        st.session_state.image_chat_history = []
    
    # Display chat history
    for i, message in enumerate(st.session_state.image_chat_history):
        if message["role"] == "user":
            with st.chat_message("user"):
                st.write(message["content"])
        else:
            with st.chat_message("assistant"):
                st.write(message["content"])
                if message.get("image_path"):
                    if message["image_path"].endswith('.svg'):
                        # Display SVG content
                        with open(message["image_path"], "r") as f:
                            svg_content = f.read()
                        st.image(svg_content, width=400)
                    else:
                        st.image(message["image_path"], width=400)
                    
                    if st.button(f"Add to Presentation", key=f"add_img_{i}"):
                        if st.session_state.current_presentation:
                            # Add to last slide
                            slide_count = len(st.session_state.current_presentation.slides)
                            if slide_count > 0:
                                st.session_state.presentation_generator.add_image_to_slide(
                                    st.session_state.current_presentation,
                                    slide_count - 1,
                                    message["image_path"]
                                )
                                st.success("Image added to presentation!")
                        else:
                            st.warning("Create a presentation first!")
    
    # Chat input
    user_prompt = st.chat_input("Describe the image you want to generate...")
    
    if user_prompt:
        # Add user message to chat
        st.session_state.image_chat_history.append({
            "role": "user",
            "content": user_prompt
        })
        
        # Generate image
        with st.spinner("Generating image..."):
            svg_content = st.session_state.image_generator.generate_svg_image(user_prompt)
            image_filename = f"generated_{len(st.session_state.generated_images)}"
            image_path = st.session_state.image_generator.save_svg_as_image(svg_content, image_filename)
            
            st.session_state.generated_images.append({
                "prompt": user_prompt,
                "path": image_path
            })
            
            # Add assistant response to chat
            st.session_state.image_chat_history.append({
                "role": "assistant",
                "content": f"Generated image for: {user_prompt}",
                "image_path": image_path
            })
        
        st.rerun()

# Current Presentation Section
if st.session_state.current_presentation:
    st.markdown("---")
    st.header("Current Presentation")
    
    prs = st.session_state.current_presentation
    slides_preview = st.session_state.presentation_generator.get_slide_preview_text(prs)
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        st.subheader(f"Presentation Overview ({len(prs.slides)} slides)")
        
        for slide_info in slides_preview:
            with st.expander(f"Slide {slide_info['slide_number']}"):
                st.text(slide_info['content'])
    
    with col2:
        st.subheader("Actions")
        
        # Download presentation
        filename = st.text_input("Filename", value="presentation.pptx")
        
        if st.button("Save & Download"):
            try:
                # Create temporary file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                    prs.save(tmp_file.name)
                    
                    # Read file for download
                    with open(tmp_file.name, 'rb') as f:
                        pptx_data = f.read()
                    
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_data,
                        file_name=filename,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
                    
                    # Clean up
                    os.unlink(tmp_file.name)
                    
            except Exception as e:
                st.error(f"Error saving presentation: {str(e)}")
        
        # Clear presentation
        if st.button("New Presentation"):
            st.session_state.current_presentation = None
            st.rerun()
        
        # Presentation stats
        st.markdown("### Statistics")
        st.metric("Total Slides", len(prs.slides))
        st.metric("Generated Images", len(st.session_state.generated_images))

# Generated Images Gallery
if st.session_state.generated_images:
    st.markdown("---")
    st.header("Generated Images Gallery")
    
    cols = st.columns(3)
    for i, img_info in enumerate(st.session_state.generated_images):
        with cols[i % 3]:
            st.markdown(f"**{img_info['prompt'][:30]}...**")
            if img_info['path'].endswith('.svg'):
                with open(img_info['path'], "r") as f:
                    svg_content = f.read()
                st.image(svg_content, width=200)
            else:
                st.image(img_info['path'], width=200)
            
            if st.button(f"Use in Slide", key=f"gallery_use_{i}"):
                if st.session_state.current_presentation:
                    slide_count = len(st.session_state.current_presentation.slides)
                    if slide_count > 0:
                        st.session_state.presentation_generator.add_image_to_slide(
                            st.session_state.current_presentation,
                            slide_count - 1,
                            img_info['path']
                        )
                        st.success("Image added!")
                else:
                    st.warning("Create a presentation first!")

# Navigation
st.markdown("---")
col1, col2 = st.columns(2)

with col1:
    if st.button("Back to Dashboard", use_container_width=True):
        st.switch_page("pages/6_Dashboard.py")

with col2:
    if st.button("Back to Analysis", use_container_width=True):
        st.switch_page("pages/5_Analysis.py")

# Sidebar help
st.sidebar.markdown("""
### Presentation Builder Guide

**Comprehensive Generation:**
- Automatically includes ALL available data
- Business problem, analysis, dashboard insights
- Latest chat conversation (most recent only)
- Sentiment analysis and predictions
- Professional templates with consistent styling

**Data Sources Included:**
- Business problem definition
- Dataset statistics and quality metrics
- Analysis results and insights
- Sentiment analysis findings
- Predictions and forecasts
- Dashboard insights
- Latest AI chat conversation
- Actionable recommendations

**AI Image Studio:**
- Chat-based image generation using Groq AI
- Professional business graphics
- Direct integration with presentation slides
- Gallery management for image reuse

**Tips:**
- Complete more workflow steps for richer presentations
- Use descriptive prompts for better AI images
- All templates maintain professional consistency
- Only latest chat conversation is included (not full history)
""")